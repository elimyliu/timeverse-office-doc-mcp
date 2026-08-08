"""模板管理 - 跨格式的统一模板管理系统。

对应方案 5.5 模板管理工具集 + 5.5.1 模板存储设计。
注册 = 复制进 templates/{format}/ + 写入注册索引 templates/registry.json。
"""

from __future__ import annotations

import json
import logging
import re
import shutil
from pathlib import Path
from typing import Any

from ..config import ServerConfig
from .error_handler import ToolError
from .path_guard import path_guard

logger = logging.getLogger("timeverse_office_doc_mcp.template")

# 占位符正则：{{variable}}, {{date:fmt}}, {{number:fmt}}, {{table:name}} 等
PLACEHOLDER_RE = re.compile(r"\{\{([^}]+)\}\}")

VALID_FORMATS = {"word", "excel", "ppt", "pdf"}


class TemplateManager:
    """跨格式模板管理器。"""

    def __init__(self, template_dir: str | None = None) -> None:
        self.template_dir = Path(template_dir or ServerConfig.TEMPLATE_DIR)
        self.registry_path = self.template_dir / "registry.json"
        self.template_dir.mkdir(parents=True, exist_ok=True)
        self._ensure_registry()

    def _ensure_registry(self) -> None:
        """确保注册索引文件存在。"""
        if not self.registry_path.exists():
            self._write_registry({})

    def _read_registry(self) -> dict[str, dict[str, Any]]:
        """读取注册索引。"""
        if not self.registry_path.exists():
            return {}
        try:
            with open(self.registry_path, encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, OSError):
            return {}

    def _write_registry(self, registry: dict[str, dict[str, Any]]) -> None:
        """写入注册索引。"""
        with open(self.registry_path, "w", encoding="utf-8") as f:
            json.dump(registry, f, ensure_ascii=False, indent=2)

    def register_template(
        self,
        name: str,
        format: str,
        file_path: str,
        description: str = "",
        placeholders: list[str] | None = None,
    ) -> dict[str, Any]:
        """注册新模板：将外部文件复制到模板库并写入索引。"""
        if format not in VALID_FORMATS:
            raise ToolError(f"不支持的格式: {format}，允许: {', '.join(sorted(VALID_FORMATS))}")

        # 校验源文件路径安全
        validated_src = path_guard.validate_path(file_path, "read")

        src = Path(validated_src)
        if not src.exists():
            raise ToolError(f"模板源文件不存在: {file_path}")

        # 复制到 templates/{format}/{name}{ext}
        format_dir = self.template_dir / format
        format_dir.mkdir(parents=True, exist_ok=True)
        dest = format_dir / f"{name}{src.suffix}"
        shutil.copy2(str(src), str(dest))

        # 自动提取占位符（如未显式提供）
        if placeholders is None:
            placeholders = self._extract_placeholders_from_file(str(dest), format)

        # 写入注册索引
        registry = self._read_registry()
        registry[name] = {
            "name": name,
            "format": format,
            "path": str(dest),
            "description": description,
            "placeholders": placeholders,
        }
        self._write_registry(registry)

        logger.info("Registered template '%s' (%s) -> %s", name, format, dest)
        return registry[name]

    def delete_template(self, template_name: str) -> None:
        """删除模板：移除索引记录并清理模板库副本（不影响用户原始文件）。"""
        registry = self._read_registry()
        entry = registry.pop(template_name, None)
        if entry is None:
            raise ToolError(f"模板不存在: {template_name}")

        # 清理模板库副本
        template_path = Path(entry["path"])
        if template_path.exists():
            template_path.unlink()

        self._write_registry(registry)
        logger.info("Deleted template '%s'", template_name)

    def list_templates(self, format: str | None = None) -> list[dict[str, Any]]:
        """列出所有可用模板，可选按格式过滤。"""
        registry = self._read_registry()
        templates = list(registry.values())
        if format:
            templates = [t for t in templates if t.get("format") == format]
        return templates

    def get_template_info(self, template_name: str) -> dict[str, Any]:
        """获取模板详情（含占位符列表与路径）。"""
        registry = self._read_registry()
        entry = registry.get(template_name)
        if entry is None:
            raise ToolError(f"模板不存在: {template_name}")
        return entry

    def resolve_template_path(self, template_name: str) -> tuple[str, str]:
        """解析模板名 -> (格式, 文件路径)。"""
        info = self.get_template_info(template_name)
        return info["format"], info["path"]

    def extract_placeholders(self, template_name: str) -> list[dict[str, str]]:
        """扫描提取模板中的占位符变量及其类型推断。"""
        info = self.get_template_info(template_name)
        return self._extract_placeholders_from_file(info["path"], info["format"])

    def _extract_placeholders_from_file(self, file_path: str, format: str) -> list[dict[str, str]]:
        """从文件中提取占位符（基础实现：读取文本并正则匹配）。"""
        try:
            text = self._read_template_text(file_path, format)
        except Exception:
            logger.warning("Failed to read template text from %s", file_path)
            return []

        seen: set[str] = set()
        placeholders: list[dict[str, str]] = []
        for match in PLACEHOLDER_RE.finditer(text):
            raw = match.group(1)
            if raw in seen:
                continue
            seen.add(raw)
            placeholders.append({"name": raw, "type": self._infer_type(raw)})
        return placeholders

    def _read_template_text(self, file_path: str, format: str) -> str:
        """按格式读取模板文本内容（用于占位符提取）。"""
        path = Path(file_path)
        if format == "word":
            from docx import Document

            doc = Document(str(path))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return "\n".join(parts)
        if format == "excel":
            from openpyxl import load_workbook

            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.extend(str(c) for c in row if c is not None)
            wb.close()
            return "\n".join(parts)
        if format == "ppt":
            from pptx import Presentation

            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return "\n".join(parts)
        if format == "pdf":
            # PDF 基础提取
            try:
                import pypdf

                reader = pypdf.PdfReader(str(path))
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except ImportError:
                return ""
        return ""

    @staticmethod
    def _infer_type(raw: str) -> str:
        """推断占位符类型。"""
        if raw.startswith(
            ("date:", "number:", "table:", "image:", "condition:", "loop:", "row:", "end:")
        ):
            return raw.split(":", 1)[0]
        return "text"


# 全局单例
template_manager = TemplateManager()
