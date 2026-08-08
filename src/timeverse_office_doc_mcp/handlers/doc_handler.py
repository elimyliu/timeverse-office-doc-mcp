"""跨格式工具处理器 - 模板管理 5 个 + Session 管理 3 个 = 8 个工具。

对应方案 5.5 模板管理工具集 + 5.6 Session 管理工具集。
"""

from __future__ import annotations

import copy
import logging
from typing import Any

from ..common.error_handler import ToolError
from ..common.file_lock import file_lock_mgr
from ..common.path_guard import path_guard
from ..common.session import session_manager
from ..common.template_mgr import template_manager
from ..common.template_utils import (
    fill_excel_variables,
    fill_ppt_slide_variables,
    fill_ppt_variables,
    fill_word_variables,
)
from ..common.validator import InputValidator

logger = logging.getLogger("timeverse_office_doc_mcp.doc")


# ==================== 5.5.1 模板注册与管理（4 个） ====================


def doc_list_templates(format: str | None = None) -> dict[str, Any]:
    """列出所有可用模板，可选按格式过滤。"""
    templates = template_manager.list_templates(format)
    return {
        "count": len(templates),
        "templates": [
            {"name": t["name"], "format": t["format"], "description": t.get("description", "")}
            for t in templates
        ],
    }


def doc_get_template_info(template_name: str) -> dict[str, Any]:
    """获取模板详情（含占位符列表与预览）。"""
    info = template_manager.get_template_info(template_name)
    return {
        "name": info["name"],
        "format": info["format"],
        "description": info.get("description", ""),
        "placeholders": info.get("placeholders", []),
        "path": info["path"],
    }


def doc_manage_template(
    action: str,
    name: str,
    format: str | None = None,
    file_path: str | None = None,
    description: str = "",
    placeholders: list[str] | None = None,
) -> dict[str, Any]:
    """模板注册与删除统一入口。

    action="register"：注册新模板到模板库（将外部文件复制到 templates/{format}/ 并写入索引）。
    action="delete"：从模板库删除模板（删除索引记录与模板库副本，不影响用户原始文件）。
    """
    InputValidator.validate_choice(action, ["register", "delete"], "action")

    if action == "register":
        result = template_manager.register_template(
            name=name,
            format=format,
            file_path=file_path,
            description=description,
            placeholders=placeholders,
        )
        return {
            "name": result["name"],
            "format": result["format"],
            "description": result.get("description", ""),
            "placeholders": result.get("placeholders", []),
            "registered": True,
        }

    # action == "delete"
    template_manager.delete_template(name)
    return {"deleted": name}


# ==================== 5.5.2 模板应用与占位符（2 个） ====================


def doc_apply_template(
    template_name: str,
    output_path: str,
    variables: dict[str, Any] | None = None,
    sections: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    """从模板创建文档并自动填充变量。

    核心工具：加载模板 -> 替换占位符 -> 保存输出。
    支持所有格式（Word/Excel/PPT/PDF）。

    sections（仅 PPT 模板生效）：按章节扩展页数。模板中需包含章节页
    （含 {{section_no}}/{{section_title}} 占位符）与内容页
    （含 {{slide_title}}/{{point1}} 占位符）原型，工具会将它们复制为
    每章节一组，插入到封面/目录与结尾页之间。每项格式：
    {
      "section_no": "01", "section_title": "公司概况",
      "slide_title": "公司概况",
      "point1": "...", "point2": "...", "point3": "...", "point4": "..."
    }
    """
    fmt, tpl_path = template_manager.resolve_template_path(template_name)
    variables = variables or {}

    if fmt == "word":
        return _apply_word_template(template_name, tpl_path, output_path, variables)
    if fmt == "excel":
        return _apply_excel_template(template_name, tpl_path, output_path, variables)
    if fmt == "ppt":
        return _apply_ppt_template(template_name, tpl_path, output_path, variables, sections)
    if fmt == "pdf":
        from .pdf_handler import pdf_create_from_template

        result = pdf_create_from_template(template_name, variables, output_path)
        return {"template": template_name, "output": result["output"], "format": "pdf"}
    raise ToolError(f"不支持的模板格式: {fmt}")


def _apply_word_template(
    template_name: str, tpl_path: str, output_path: str, variables: dict[str, Any]
) -> dict[str, Any]:
    """应用 Word 模板。"""
    from docx import Document

    validated_out = path_guard.validate_path(output_path, "write")
    doc = Document(tpl_path)
    replaced = fill_word_variables(doc, variables)
    file_lock_mgr.acquire(validated_out)
    try:
        doc.save(validated_out)
    finally:
        file_lock_mgr.release(validated_out)
    return {
        "template": template_name,
        "output": validated_out,
        "format": "word",
        "variables_replaced": replaced,
    }


def _apply_excel_template(
    template_name: str, tpl_path: str, output_path: str, variables: dict[str, Any]
) -> dict[str, Any]:
    """应用 Excel 模板。"""
    from openpyxl import load_workbook

    validated_out = path_guard.validate_path(output_path, "write")
    wb = load_workbook(tpl_path)
    replaced = fill_excel_variables(wb, variables)
    file_lock_mgr.acquire(validated_out)
    try:
        wb.save(validated_out)
    finally:
        file_lock_mgr.release(validated_out)
    return {
        "template": template_name,
        "output": validated_out,
        "format": "excel",
        "variables_replaced": replaced,
    }


def _apply_ppt_template(
    template_name: str,
    tpl_path: str,
    output_path: str,
    variables: dict[str, Any],
    sections: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    """应用 PPT 模板。

    sections 提供时先按章节扩展页数，再填充变量。
    """
    from pptx import Presentation

    validated_out = path_guard.validate_path(output_path, "write")
    prs = Presentation(tpl_path)
    if sections:
        replaced = _expand_ppt_sections(prs, sections, variables)
    else:
        replaced = fill_ppt_variables(prs, variables)
    file_lock_mgr.acquire(validated_out)
    try:
        prs.save(validated_out)
    finally:
        file_lock_mgr.release(validated_out)
    return {
        "template": template_name,
        "output": validated_out,
        "format": "ppt",
        "variables_replaced": replaced,
        "slide_count": len(prs.slides),
    }


# ==================== PPT 多章节扩页辅助 ====================


def _find_proto_idx(slides: list[Any], hints: tuple[str, ...]) -> int | None:
    """查找第一个包含任一指定占位符特征的页面索引。"""
    for idx, slide in enumerate(slides):
        if any(hint in _collect_slide_text(slide) for hint in hints):
            return idx
    return None


def _collect_slide_text(slide: Any) -> str:
    """收集幻灯片中文本框与表格的全部文本。"""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "".join(parts)


def _append_cloned_slide(prs: Any, layout: Any, shapes: list[Any]) -> Any:
    """追加一页，将深拷贝的形状 XML 粘贴进去，返回新幻灯片对象。"""
    new_slide = prs.slides.add_slide(layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for el in shapes:
        new_slide.shapes._spTree.insert_element_before(copy.deepcopy(el), "p:extLst")
    return new_slide


def _reorder_ppt_slides(prs: Any, final_order: list[Any]) -> None:
    """按给定幻灯片对象顺序重排文档。"""
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    slides = list(prs.slides)
    obj_to_el = {id(slides[i]): sld_ids[i] for i in range(len(slides))}
    for el in sld_ids:
        sld_id_lst.remove(el)
    for slide in final_order:
        sld_id_lst.append(obj_to_el[id(slide)])


def _resequence_ppt_page_numbers(prs: Any) -> None:
    """按最终页序重写页脚页码。

    扩页时克隆页会复制模板中的静态页码（如 "03"/"04"），导致页码重复。
    本函数识别「位于幻灯片底部区域的纯数字文本框」作为页码，并按 1-based
    页序重写，保留原数字位数格式（02 / 002）。
    """
    slide_height = prs.slide_height
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text.isdigit() or len(text) > 3:
                continue
            # 页码判定：文本框顶部位于幻灯片下部（bottom 15% 区域）
            if shape.top is None or shape.top < slide_height * 0.85:
                continue
            new_text = str(idx).zfill(len(text))
            if new_text != text:
                para = shape.text_frame.paragraphs[0]
                runs = para.runs
                if runs:
                    runs[0].text = new_text
                    for run in runs[1:]:
                        run.text = ""


def _expand_ppt_sections(
    prs: Any, sections: list[dict[str, Any]], variables: dict[str, Any]
) -> int:
    """按章节扩展 PPT：以章节页/内容页原型为模板生成每章节一组页面。

    第 1 章直接复用模板中的原型页，第 2..N 章克隆副本，随后重排为
    封面/目录 + 章节页 xN + 内容页 xN + 结尾页，并填充变量。
    返回替换的占位符数量。
    """
    slides = list(prs.slides)
    sec_idx = _find_proto_idx(slides, ("section_no", "section_title"))
    con_idx = _find_proto_idx(slides, ("slide_title", "point1"))

    if sec_idx is None and con_idx is None:
        # 模板不含章节/内容原型页，退化为仅做全局变量填充
        return fill_ppt_variables(prs, variables)

    proto_idx = sorted({i for i in (sec_idx, con_idx) if i is not None})
    proto_max = proto_idx[-1]

    # 头部页（封面/目录等）与尾部页（结尾页等）保持原位
    head = [slides[i] for i in range(proto_max + 1) if i not in proto_idx]
    tail = slides[proto_max + 1 :]

    # 快照原型页形状（用于生成副本）
    def _snapshot(idx: int | None) -> tuple[Any, list[Any]] | None:
        if idx is None:
            return None
        slide = slides[idx]
        return slide.slide_layout, [copy.deepcopy(sh._element) for sh in slide.shapes]

    sec_snap = _snapshot(sec_idx)
    con_snap = _snapshot(con_idx)
    per_section = (1 if sec_snap is not None else 0) + (1 if con_snap is not None else 0)

    # 追加第 2..N 章的副本（第 1 章直接复用原型页）
    clones: list[Any] = []
    for sec in sections[1:]:
        if sec_snap is not None:
            clones.append(_append_cloned_slide(prs, sec_snap[0], sec_snap[1]))
        if con_snap is not None:
            clones.append(_append_cloned_slide(prs, con_snap[0], con_snap[1]))

    # 目标顺序：头部 + 第1章(原型) + 第2..N章(副本) + 尾部
    protos = [slides[i] for i in proto_idx]
    _reorder_ppt_slides(prs, head + protos + clones + tail)

    # 先全局填充头部/尾部页（封面、目录、结尾），再逐章节填充新增页。
    # 注意：不能全文档填充，否则章节占位符会被缺失键替换为空串。
    replaced = 0
    for slide in head + tail:
        replaced += fill_ppt_slide_variables(slide, variables)
    base = len(head)
    for i, sec in enumerate(sections):
        pos = base + i * per_section
        if sec_snap is not None:
            replaced += fill_ppt_slide_variables(
                prs.slides[pos],
                {
                    "section_no": str(sec.get("section_no", "")),
                    "section_title": str(sec.get("section_title", "")),
                },
            )
        if con_snap is not None:
            con_vars: dict[str, Any] = {"slide_title": str(sec.get("slide_title", ""))}
            for j in range(1, 5):
                con_vars[f"point{j}"] = str(sec.get(f"point{j}", ""))
            replaced += fill_ppt_slide_variables(
                prs.slides[pos + (1 if sec_snap is not None else 0)], con_vars
            )
    # 克隆页会复制模板静态页码，按最终页序重写页脚页码
    _resequence_ppt_page_numbers(prs)
    return replaced


def doc_extract_placeholders(template_name: str, format: str | None = None) -> dict[str, Any]:
    """扫描提取模板中的占位符变量。"""
    placeholders = template_manager.extract_placeholders(template_name)
    return {
        "template_name": template_name,
        "count": len(placeholders),
        "placeholders": placeholders,
    }


# ==================== 5.6 Session 管理工具集（4 个） ====================


def doc_open_session(filename: str, format: str) -> dict[str, Any]:
    """打开文档到内存 Session。

    返回 session_id，后续编辑工具可通过 session_id 操作内存中的文档。
    """
    validated = path_guard.validate_path(filename, "read")

    if format == "word":
        from docx import Document

        doc = Document(validated)
    elif format == "excel":
        from openpyxl import load_workbook

        doc = load_workbook(validated)
    elif format == "ppt":
        from pptx import Presentation

        doc = Presentation(validated)
    elif format == "pdf":
        from pypdf import PdfReader

        doc = PdfReader(validated)
    else:
        raise ToolError(f"不支持的格式: {format}")

    session_id = session_manager.open_session(validated, format, doc)
    return {
        "session_id": session_id,
        "filename": validated,
        "format": format,
    }


def doc_save_session(session_id: str, output_path: str | None = None) -> dict[str, Any]:
    """保存 Session 到磁盘。

    不指定 output_path 则保存回原路径。
    """
    session = session_manager.get_session(session_id)

    if output_path:
        validated_out = path_guard.validate_path(output_path, "write")
    else:
        validated_out = session.filename

    # 按格式保存
    doc = session.document
    if session.format == "word" or session.format == "excel" or session.format == "ppt":
        file_lock_mgr.acquire(validated_out)
        try:
            doc.save(validated_out)
        finally:
            file_lock_mgr.release(validated_out)
    elif session.format == "pdf":
        # PDF 的 Session 保存逻辑（如果有 PdfWriter 则写入）
        if hasattr(doc, "write"):
            file_lock_mgr.acquire(validated_out)
            try:
                with open(validated_out, "wb") as f:
                    doc.write(f)
            finally:
                file_lock_mgr.release(validated_out)

    session_manager.save_session(session_id, validated_out)
    return {"session_id": session_id, "saved_to": validated_out}


def doc_close_session(
    session_id: str,
    save: bool = False,
    output_path: str | None = None,
) -> dict[str, Any]:
    """关闭 Session。

    save=True 时先保存再关闭。
    output_path 指定时保存到该路径（仅在 save=True 时生效）。
    """
    if save:
        doc_save_session(session_id, output_path)
    session_manager.close_session(session_id)
    return {"session_id": session_id, "closed": True, "saved": save}


def doc_list_sessions() -> dict[str, Any]:
    """列出所有活跃 Session。"""
    sessions = session_manager.list_sessions()
    return {
        "count": len(sessions),
        "sessions": sessions,
    }
