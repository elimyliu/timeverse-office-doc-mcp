"""PowerPoint 演示文稿处理器 - 19 个工具。

对应方案 5.3 PowerPoint 工具集。
使用 python-pptx 实现，支持 Session 内存编辑模式。
"""

from __future__ import annotations

import copy
import logging
import re
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

from ..common.error_handler import ToolError
from ..common.file_lock import file_lock_mgr
from ..common.path_guard import path_guard
from ..common.session import session_manager
from ..common.template_mgr import template_manager
from ..common.template_utils import fill_ppt_variables
from ..common.validator import InputValidator

logger = logging.getLogger("timeverse_office_doc_mcp.ppt")

# EMU 与英寸的换算常数（914400 EMU = 1 英寸）
_EMU_PER_INCH = 914400

# 图表类型映射
_CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}

# 形状类型映射（textbox 特殊处理，走 add_textbox）
_SHAPE_TYPE_MAP: dict[str, MSO_SHAPE | None] = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "textbox": None,
}

# 主题配色方案（修改 a:clrScheme 的 dk2 与 accent1-6）
_THEME_COLORS: dict[str, dict[str, str]] = {
    "blue": {
        "dk2": "1F3864",
        "accent1": "4472C4",
        "accent2": "5B9BD5",
        "accent3": "70AD47",
        "accent4": "FFC000",
        "accent5": "4BACC6",
        "accent6": "8064A2",
    },
    "green": {
        "dk2": "1F4E2C",
        "accent1": "2E7D32",
        "accent2": "4CAF50",
        "accent3": "81C784",
        "accent4": "A5D6A7",
        "accent5": "66BB6A",
        "accent6": "388E3C",
    },
    "orange": {
        "dk2": "8B4500",
        "accent1": "ED7D31",
        "accent2": "F4B183",
        "accent3": "C55A11",
        "accent4": "FFC000",
        "accent5": "FFD966",
        "accent6": "BF9000",
    },
    "dark": {
        "dk2": "0D0D0D",
        "accent1": "262626",
        "accent2": "404040",
        "accent3": "595959",
        "accent4": "808080",
        "accent5": "A6A6A6",
        "accent6": "BFBFBF",
    },
}


# ==================== 辅助函数 ====================


def _get_presentation(filename: str, session_id: str | None = None) -> Presentation:
    """获取 Presentation 对象：Session 模式从内存取，否则从磁盘打开。"""
    if session_id:
        return session_manager.get_document(session_id, "ppt")
    validated = path_guard.validate_path(filename, "read")
    return Presentation(validated)


def _save_presentation(prs: Presentation, filename: str, session_id: str | None = None) -> None:
    """保存演示文稿：Session 模式仅标记修改，否则写入磁盘。"""
    if session_id:
        session_manager.mark_modified(session_id)
    else:
        validated = path_guard.validate_path(filename, "write")
        file_lock_mgr.acquire(validated)
        try:
            prs.save(validated)
        finally:
            file_lock_mgr.release(validated)


def _get_slide(prs: Presentation, slide_idx: int) -> Any:
    """按索引获取幻灯片，索引越界时抛出 ToolError。"""
    slide_count = len(prs.slides)
    if slide_idx < 0 or slide_idx >= slide_count:
        raise ToolError(f"幻灯片索引超出范围: {slide_idx}（共 {slide_count} 张）")
    return prs.slides[slide_idx]


def _emu_to_inches(emu: int | None) -> float | None:
    """将 EMU 转为英寸，None 则返回 None。"""
    if emu is None:
        return None
    return round(emu / _EMU_PER_INCH, 2)


def _get_theme_part(prs: Presentation) -> Any:
    """获取演示文稿主题部件（theme1.xml），找不到则抛出 ToolError。"""
    master = prs.slide_master
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            return rel.target_part
    raise ToolError("未找到演示文稿主题部件")


def _set_scheme_color(clr_scheme: Any, name: str, hex_val: str) -> None:
    """设置颜色方案中指定名称的颜色（统一写为 a:srgbClr）。"""
    elem = clr_scheme.find(qn(f"a:{name}"))
    if elem is None:
        return
    srgb = elem.find(qn("a:srgbClr"))
    sys_clr = elem.find(qn("a:sysClr"))
    if srgb is not None:
        srgb.set("val", hex_val)
    elif sys_clr is not None:
        elem.remove(sys_clr)
        elem.append(parse_xml(f'<a:srgbClr {nsdecls("a")} val="{hex_val}"/>'))
    else:
        elem.append(parse_xml(f'<a:srgbClr {nsdecls("a")} val="{hex_val}"/>'))


def _apply_theme_colors(prs: Presentation, theme_name: str) -> dict[str, str]:
    """将主题配色方案写入主题部件的颜色方案，返回应用的配色字典。"""
    colors = _THEME_COLORS.get(theme_name)
    if colors is None:
        raise ToolError(f"不支持的主题: {theme_name}，可选: {', '.join(sorted(_THEME_COLORS))}")
    theme_part = _get_theme_part(prs)
    root = etree.fromstring(theme_part.blob)
    clr_scheme = root.find(qn("a:themeElements")).find(qn("a:clrScheme"))
    for name, hex_val in colors.items():
        _set_scheme_color(clr_scheme, name, hex_val)
    theme_part._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
    return colors


def _normalize_chart_series(
    data: dict[str, Any],
) -> tuple[list[str], list[tuple[str, list[float]]]]:
    """规范化图表数据，返回 (categories, [(series_name, values)])。"""
    categories = data.get("categories", [])
    raw_series = data.get("series", [])
    series_list: list[tuple[str, list[float]]] = []
    if isinstance(raw_series, dict):
        for name, values in raw_series.items():
            series_list.append((str(name), [float(v) for v in values]))
    elif isinstance(raw_series, list):
        for item in raw_series:
            if isinstance(item, dict):
                name = str(item.get("name", "Series"))
                values = [float(v) for v in item.get("values", [])]
            elif isinstance(item, (list, tuple)) and len(item) == 2:
                name = str(item[0])
                values = [float(v) for v in item[1]]
            else:
                continue
            series_list.append((name, values))
    if not categories or not series_list:
        raise ToolError("图表数据需包含 categories 与至少一组 series")
    return categories, series_list


def _shape_type_name(shape: Any) -> str:
    """获取形状类型的可读名称。"""
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        return f"placeholder({ph_type})" if ph_type else "placeholder"
    return shape.shape_type.name if shape.shape_type else "unknown"


# ==================== 5.3.1 演示文稿管理（7 个） ====================


def ppt_create_presentation(
    filename: str,
    template: str | None = None,
    variables: dict[str, Any] | None = None,
    session_id: str | None = None,
) -> dict[str, Any]:
    """创建新演示文稿（支持模板与变量填充）。"""
    InputValidator.validate_filename(filename)
    validated = path_guard.validate_path(filename, "write")

    if template:
        fmt, tpl_path = template_manager.resolve_template_path(template)
        if fmt != "ppt":
            raise ToolError(f"模板 '{template}' 是 {fmt} 格式，不是 ppt")
        prs = Presentation(tpl_path)
        replaced = fill_ppt_variables(prs, variables) if variables else 0
    else:
        prs = Presentation()
        replaced = 0

    actual_sid: str | None = None
    if session_id:
        actual_sid = session_manager.open_session(validated, "ppt", prs)
    else:
        file_lock_mgr.acquire(validated)
        try:
            prs.save(validated)
        finally:
            file_lock_mgr.release(validated)

    return {
        "filename": validated,
        "template": template,
        "variables_replaced": replaced,
        "slide_count": len(prs.slides),
        "session_id": actual_sid,
    }


def ppt_get_info(filename: str, session_id: str | None = None) -> dict[str, Any]:
    """获取演示文稿元信息（尺寸、幻灯片数、模板等）。"""
    prs = _get_presentation(filename, session_id)
    props = prs.core_properties
    return {
        "filename": filename,
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
        "slide_count": len(prs.slides),
        "slide_width": _emu_to_inches(prs.slide_width),
        "slide_height": _emu_to_inches(prs.slide_height),
        "layout_count": len(prs.slide_layouts),
    }


def ppt_list_slides(filename: str, session_id: str | None = None) -> dict[str, Any]:
    """列出所有幻灯片概览。"""
    prs = _get_presentation(filename, session_id)
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        shape_count = len(slide.shapes)
        text_preview = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_preview = shape.text_frame.text.strip()[:50]
                break
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text[:80]
        slides.append(
            {
                "slide_idx": idx,
                "layout": slide.slide_layout.name,
                "shape_count": shape_count,
                "text_preview": text_preview,
                "has_notes": slide.has_notes_slide,
                "notes_preview": notes_text,
            }
        )
    return {"filename": filename, "slide_count": len(slides), "slides": slides}


def ppt_add_slide(
    filename: str,
    layout: int = 6,
    title: str | None = None,
    session_id: str | None = None,
) -> dict[str, Any]:
    """添加幻灯片（layout: 0-8 对应 python-pptx slide_layouts）。"""
    prs = _get_presentation(filename, session_id)
    if layout < 0 or layout >= len(prs.slide_layouts):
        raise ToolError(f"布局索引超出范围: {layout}（可用 0-{len(prs.slide_layouts) - 1}）")
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    if title:
        for shape in slide.placeholders:
            if (
                shape.placeholder_format.type is not None
                and "title" in str(shape.placeholder_format.type).lower()
            ):
                shape.text = title
                break
    new_idx = len(prs.slides) - 1
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": new_idx,
        "layout": layout,
        "layout_name": prs.slide_layouts[layout].name,
        "title": title,
    }


def ppt_delete_slide(
    filename: str,
    slide_idx: int,
    session_id: str | None = None,
) -> dict[str, Any]:
    """删除指定幻灯片。"""
    prs = _get_presentation(filename, session_id)
    slide_count = len(prs.slides)
    if slide_idx < 0 or slide_idx >= slide_count:
        raise ToolError(f"幻灯片索引超出范围: {slide_idx}（共 {slide_count} 张）")
    sld_id_lst = prs.slides._sldIdLst
    sld_id = list(sld_id_lst)[slide_idx]
    r_id = sld_id.get(qn("r:id"))
    sld_id_lst.remove(sld_id)
    if r_id:
        prs.part.drop_rel(r_id)
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "deleted_slide_idx": slide_idx,
        "remaining_slides": len(prs.slides),
    }


def ppt_move_slide(
    filename: str,
    slide_idx: int,
    new_idx: int,
    session_id: str | None = None,
) -> dict[str, Any]:
    """移动幻灯片到新位置。"""
    prs = _get_presentation(filename, session_id)
    slide_count = len(prs.slides)
    if slide_idx < 0 or slide_idx >= slide_count:
        raise ToolError(f"幻灯片索引超出范围: {slide_idx}（共 {slide_count} 张）")
    if new_idx < 0 or new_idx >= slide_count:
        raise ToolError(f"目标索引超出范围: {new_idx}（共 {slide_count} 张）")
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    target = sld_ids[slide_idx]
    sld_id_lst.remove(target)
    sld_id_lst.insert(new_idx, target)
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "from_idx": slide_idx,
        "to_idx": new_idx,
        "slide_count": len(prs.slides),
    }


def ppt_copy_slide(
    filename: str,
    slide_idx: int,
    session_id: str | None = None,
) -> dict[str, Any]:
    """复制幻灯片（副本插入到原幻灯片之后）。"""
    prs = _get_presentation(filename, session_id)
    slide_count = len(prs.slides)
    if slide_idx < 0 or slide_idx >= slide_count:
        raise ToolError(f"幻灯片索引超出范围: {slide_idx}（共 {slide_count} 张）")
    src_slide = prs.slides[slide_idx]
    src_layout = src_slide.slide_layout

    # 使用源版式创建新幻灯片，清空默认占位形状后深拷贝源形状
    new_slide = prs.slides.add_slide(src_layout)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    for shape in src_slide.shapes:
        new_elem = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_elem, "p:extLst")
    # 复制部件关系（图片、图表等），排除备注页与版式（add_slide 已建立版式关系）
    for rel in src_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    # 将副本从末尾移动到原幻灯片之后
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    copied = sld_ids[-1]
    sld_id_lst.remove(copied)
    sld_id_lst.insert(slide_idx + 1, copied)

    new_idx = slide_idx + 1
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "source_slide_idx": slide_idx,
        "new_slide_idx": new_idx,
        "slide_count": len(prs.slides),
    }


# ==================== 5.3.2 内容编辑（8 个） ====================


def ppt_add_text(
    filename: str,
    slide_idx: int,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 8.0,
    height: float = 1.5,
    font_size: int = 18,
    bold: bool = False,
    session_id: str | None = None,
) -> dict[str, Any]:
    """添加文本框（left/top/width/height 单位为英寸）。"""
    InputValidator.validate_text_length(text)
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = textbox.text_frame
    tf.text = text
    if tf.paragraphs:
        for run in tf.paragraphs[0].runs:
            if font_size:
                run.font.size = Pt(font_size)
            if bold:
                run.font.bold = True
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "text": text,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "font_size": font_size,
        "bold": bold,
    }


def ppt_add_image(
    filename: str,
    slide_idx: int,
    image_path: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float | None = None,
    height: float | None = None,
    session_id: str | None = None,
) -> dict[str, Any]:
    """插入图片到幻灯片。"""
    img = path_guard.validate_path(image_path, "read")
    if not Path(img).exists():
        raise ToolError(f"图片文件不存在: {image_path}")
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    kwargs: dict[str, Any] = {"left": Inches(left), "top": Inches(top)}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    picture = slide.shapes.add_picture(img, **kwargs)
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "image_path": img,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "shape_name": picture.name,
    }


def ppt_add_table(
    filename: str,
    slide_idx: int,
    rows: int,
    cols: int,
    data: list[list[str]] | None = None,
    left: float = 1.0,
    top: float = 2.0,
    width: float = 8.0,
    height: float = 3.0,
    session_id: str | None = None,
) -> dict[str, Any]:
    """添加表格到幻灯片。"""
    if rows < 1:
        raise ToolError(f"行数必须大于 0，得到: {rows}")
    if cols < 1:
        raise ToolError(f"列数必须大于 0，得到: {cols}")
    if data:
        InputValidator.validate_table_data(data)
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table
    if data:
        for i, row_data in enumerate(data):
            if i >= rows:
                break
            for j, cell_text in enumerate(row_data):
                if j >= cols:
                    break
                table.cell(i, j).text = str(cell_text)
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "rows": rows,
        "cols": cols,
        "shape_name": table_shape.name,
    }


def ppt_add_chart(
    filename: str,
    slide_idx: int,
    chart_type: str,
    data: dict[str, Any],
    title: str = "",
    session_id: str | None = None,
) -> dict[str, Any]:
    """添加图表（chart_type: bar/line/pie；data 含 categories 与 series）。"""
    InputValidator.validate_choice(chart_type, ["bar", "line", "pie"], "chart_type")
    if not isinstance(data, dict):
        raise ToolError("图表数据必须是包含 categories 与 series 的字典")
    categories, series_list = _normalize_chart_series(data)
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_list:
        chart_data.add_series(name, values)

    xl_chart_type = _CHART_TYPE_MAP[chart_type]
    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(1.0), Inches(2.0), Inches(8.0), Inches(4.5), chart_data
    )
    chart = chart_shape.chart
    if title:
        chart.chart_title.text_frame.text = title
        chart.has_title = True
    else:
        chart.has_title = False
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "chart_type": chart_type,
        "title": title,
        "categories": categories,
        "series_count": len(series_list),
        "shape_name": chart_shape.name,
    }


def ppt_add_shape(
    filename: str,
    slide_idx: int,
    shape_type: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 3.0,
    height: float = 2.0,
    session_id: str | None = None,
) -> dict[str, Any]:
    """添加形状（shape_type: rectangle/rounded_rectangle/oval/triangle/textbox）。"""
    InputValidator.validate_choice(
        shape_type,
        list(_SHAPE_TYPE_MAP),
        "shape_type",
    )
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    mso_type = _SHAPE_TYPE_MAP[shape_type]
    if mso_type is None:
        # textbox 走 add_textbox
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        shape = slide.shapes.add_shape(
            mso_type, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "shape_type": shape_type,
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "shape_name": shape.name,
    }


def ppt_set_background(
    filename: str,
    slide_idx: int,
    color: str | None = None,
    session_id: str | None = None,
) -> dict[str, Any]:
    """设置幻灯片背景色（color: 十六进制如 FF0000）。"""
    if color is None:
        raise ToolError("必须提供 color 参数（十六进制颜色，如 FF0000）")
    if not re.match(r"^[0-9A-Fa-f]{6}$", color):
        raise ToolError(f"无效的十六进制颜色: {color}（应为 6 位，如 FF0000）")
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string(color)
    _save_presentation(prs, filename, session_id)
    return {"filename": filename, "slide_idx": slide_idx, "background_color": color}


def ppt_apply_theme(
    filename: str,
    theme_name: str,
    session_id: str | None = None,
) -> dict[str, Any]:
    """应用主题配色（theme_name: blue/green/orange/dark；修改颜色方案）。"""
    InputValidator.validate_choice(theme_name, list(_THEME_COLORS), "theme_name")
    prs = _get_presentation(filename, session_id)
    applied = _apply_theme_colors(prs, theme_name)
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "theme_name": theme_name,
        "colors_applied": applied,
    }


def ppt_set_slide_notes(
    filename: str,
    slide_idx: int,
    notes_text: str,
    session_id: str | None = None,
) -> dict[str, Any]:
    """设置演讲者备注。"""
    InputValidator.validate_text_length(notes_text)
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    slide.notes_slide.notes_text_frame.text = notes_text
    _save_presentation(prs, filename, session_id)
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "notes_length": len(notes_text),
    }


# ==================== 5.3.3 分析工具（4 个） ====================


def ppt_extract_text(filename: str, session_id: str | None = None) -> dict[str, Any]:
    """提取所有幻灯片文本。"""
    prs = _get_presentation(filename, session_id)
    slides_text: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        shape_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                shape_texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            shape_texts.append(cell.text)
        slides_text.append(
            {
                "slide_idx": idx,
                "texts": shape_texts,
                "text": "\n".join(shape_texts),
            }
        )
    total_chars = sum(len(s["text"]) for s in slides_text)
    return {
        "filename": filename,
        "slide_count": len(prs.slides),
        "total_chars": total_chars,
        "slides": slides_text,
    }


def ppt_get_slide_notes(
    filename: str,
    slide_idx: int,
    session_id: str | None = None,
) -> dict[str, Any]:
    """获取指定幻灯片的演讲者备注。"""
    prs = _get_presentation(filename, session_id)
    slide = _get_slide(prs, slide_idx)
    notes_text = ""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text
    return {
        "filename": filename,
        "slide_idx": slide_idx,
        "notes": notes_text,
        "has_notes": slide.has_notes_slide,
    }


def ppt_analyze_structure(filename: str, session_id: str | None = None) -> dict[str, Any]:
    """分析结构（幻灯片分布、元素统计、布局分析）。"""
    structure = ppt_get_structure(filename, session_id)

    type_dist: dict[str, int] = {}
    layout_dist: dict[str, int] = {}
    slide_details: list[dict[str, Any]] = []
    total_shapes = 0
    total_chars = 0
    notes_count = 0

    for slide in structure["slides"]:
        layout_name = slide["layout"]
        layout_dist[layout_name] = layout_dist.get(layout_name, 0) + 1
        shape_types: dict[str, int] = {}
        slide_chars = 0
        for shape in slide["shapes"]:
            type_name = shape["type"]
            type_dist[type_name] = type_dist.get(type_name, 0) + 1
            shape_types[type_name] = shape_types.get(type_name, 0) + 1
            if "text" in shape:
                slide_chars += len(shape["text"])
        total_shapes += slide["shape_count"]
        total_chars += slide_chars
        if slide["notes"].strip():
            notes_count += 1
        slide_details.append(
            {
                "slide_idx": slide["slide_idx"],
                "layout": layout_name,
                "shape_count": slide["shape_count"],
                "shape_types": shape_types,
                "text_chars": slide_chars,
            }
        )

    slide_count = structure["slide_count"]
    avg_shapes = round(total_shapes / slide_count, 2) if slide_count else 0.0
    return {
        "filename": filename,
        "slide_count": slide_count,
        "slide_width": structure["slide_width"],
        "slide_height": structure["slide_height"],
        "total_shapes": total_shapes,
        "avg_shapes_per_slide": avg_shapes,
        "total_text_chars": total_chars,
        "slides_with_notes": notes_count,
        "shape_type_distribution": type_dist,
        "layout_distribution": layout_dist,
        "slide_details": slide_details,
    }


def ppt_get_structure(filename: str, session_id: str | None = None) -> dict[str, Any]:
    """获取完整结构树。"""
    prs = _get_presentation(filename, session_id)
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        shapes: list[dict[str, Any]] = []
        for shape in slide.shapes:
            shape_info: dict[str, Any] = {
                "name": shape.name,
                "type": _shape_type_name(shape),
                "left": _emu_to_inches(shape.left),
                "top": _emu_to_inches(shape.top),
                "width": _emu_to_inches(shape.width),
                "height": _emu_to_inches(shape.height),
            }
            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text
            if shape.has_table:
                shape_info["table_rows"] = len(shape.table.rows)
                shape_info["table_cols"] = len(shape.table.columns)
            if shape.has_chart:
                shape_info["chart_type"] = str(shape.chart.chart_type)
            if shape.is_placeholder:
                shape_info["placeholder_type"] = str(shape.placeholder_format.type)
            shapes.append(shape_info)
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
        slides.append(
            {
                "slide_idx": idx,
                "layout": slide.slide_layout.name,
                "shape_count": len(slide.shapes),
                "shapes": shapes,
                "notes": notes_text,
            }
        )
    return {
        "filename": filename,
        "slide_width": _emu_to_inches(prs.slide_width),
        "slide_height": _emu_to_inches(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": slides,
    }
