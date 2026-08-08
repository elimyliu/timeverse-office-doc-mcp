"""测试跨格式 doc handler 工具（模板管理 + Session 管理）。"""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

from timeverse_office_doc_mcp.common.error_handler import ToolError
from timeverse_office_doc_mcp.handlers import doc_handler, word_handler


def _build_ppt_section_template(path: Path) -> None:
    """构造 5 页结构模板：封面/目录/章节页原型/内容页原型/结尾页。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    def add_textbox(slide, text: str) -> None:
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        tb.text_frame.text = text

    s = prs.slides.add_slide(blank)
    add_textbox(s, "{{title}}")
    s = prs.slides.add_slide(blank)
    add_textbox(s, "{{item1}} / {{item2}} / {{item3}} / {{item4}}")
    s = prs.slides.add_slide(blank)
    add_textbox(s, "PART {{section_no}} {{section_title}}")
    s = prs.slides.add_slide(blank)
    add_textbox(s, "{{slide_title}}\n{{point1}}\n{{point2}}\n{{point3}}\n{{point4}}")
    s = prs.slides.add_slide(blank)
    add_textbox(s, "{{contact}}")
    prs.save(str(path))


@pytest.fixture
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """创建临时工作空间并更新 PathGuard 白名单。"""
    ws = tmp_path / "workspace"
    ws.mkdir()
    from timeverse_office_doc_mcp.common.path_guard import path_guard

    monkeypatch.setattr(path_guard, "allowed_dirs", [tmp_path.resolve()])
    # 同时更新模板管理器的目录
    tpl_dir = tmp_path / "templates"
    tpl_dir.mkdir(exist_ok=True)
    from timeverse_office_doc_mcp.common.template_mgr import template_manager

    monkeypatch.setattr(template_manager, "template_dir", tpl_dir)
    monkeypatch.setattr(template_manager, "registry_path", tpl_dir / "registry.json")
    return ws


@pytest.fixture
def template_docx(workspace: Path) -> str:
    """创建一个带占位符的 Word 模板文件。"""
    path = str(workspace / "template.docx")
    doc = Document()
    doc.add_heading("{{title}}", level=1)
    doc.add_paragraph("作者: {{author}}")
    doc.add_paragraph("日期: {{date}}")
    doc.save(path)
    return path


class TestDocRegisterAndDelete:
    def test_register_and_list(self, template_docx: str) -> None:
        # 注册模板
        result = doc_handler.doc_manage_template(
            action="register",
            name="test_tpl",
            format="word",
            file_path=template_docx,
            description="测试模板",
        )
        assert result["registered"] is True
        assert result["name"] == "test_tpl"

        # 列出模板
        listing = doc_handler.doc_list_templates(format="word")
        assert listing["count"] >= 1
        names = [t["name"] for t in listing["templates"]]
        assert "test_tpl" in names

    def test_get_template_info(self, template_docx: str) -> None:
        doc_handler.doc_manage_template(
            action="register",
            name="info_tpl",
            format="word",
            file_path=template_docx,
            description="信息测试",
        )
        info = doc_handler.doc_get_template_info("info_tpl")
        assert info["name"] == "info_tpl"
        assert info["format"] == "word"
        assert len(info["placeholders"]) >= 2  # title, author, date

    def test_delete_template(self, template_docx: str) -> None:
        doc_handler.doc_manage_template(
            action="register",
            name="del_tpl",
            format="word",
            file_path=template_docx,
        )
        doc_handler.doc_manage_template(action="delete", name="del_tpl")
        listing = doc_handler.doc_list_templates()
        names = [t["name"] for t in listing["templates"]]
        assert "del_tpl" not in names


class TestDocApplyTemplate:
    def test_apply_word_template(self, template_docx: str, workspace: Path) -> None:
        # 注册模板
        doc_handler.doc_manage_template(
            action="register",
            name="apply_tpl",
            format="word",
            file_path=template_docx,
        )

        # 应用模板
        output = str(workspace / "output.docx")
        result = doc_handler.doc_apply_template(
            template_name="apply_tpl",
            output_path=output,
            variables={"title": "测试报告", "author": "张三", "date": "2026-08-06"},
        )
        assert result["format"] == "word"
        assert result["variables_replaced"] >= 2
        assert Path(output).exists()

        # 验证内容
        doc = Document(output)
        text = "\n".join(p.text for p in doc.paragraphs)
        assert "测试报告" in text
        assert "张三" in text
        assert "{{" not in text

    def test_apply_ppt_template_returns_name(self, workspace: Path) -> None:
        """doc_apply_template 对 PPT 模板应返回真实模板名（回归：曾硬编码为 'ppt'）。"""
        tpl = str(workspace / "tpl.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = "标题：{{title}}"
        prs.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="ppt_tpl", format="ppt", file_path=tpl
        )

        output = str(workspace / "out.pptx")
        result = doc_handler.doc_apply_template("ppt_tpl", output, {"title": "公司介绍"})
        assert result["template"] == "ppt_tpl"

        prs2 = Presentation(output)
        text = prs2.slides[0].shapes[0].text_frame.text
        assert "公司介绍" in text
        assert "{{" not in text

    def test_apply_excel_template_returns_name(self, workspace: Path) -> None:
        """doc_apply_template 对 Excel 模板应返回真实模板名（回归：曾硬编码为 'excel'）。"""
        tpl = str(workspace / "tpl.xlsx")
        wb = Workbook()
        wb.active["A1"] = "{{name}}"
        wb.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="excel_tpl", format="excel", file_path=tpl
        )

        output = str(workspace / "out.xlsx")
        result = doc_handler.doc_apply_template("excel_tpl", output, {"name": "银杉"})
        assert result["template"] == "excel_tpl"

        wb2 = load_workbook(output)
        assert wb2.active["A1"].value == "银杉"

    def test_apply_word_template_split_runs(self, workspace: Path) -> None:
        """占位符被拆到多个 run 时应能替换（回归：原先只按单个 run 替换）。"""
        tpl = str(workspace / "split.docx")
        doc = Document()
        p = doc.add_paragraph()
        p.add_run("报告：{{ti")
        p.add_run("tle}}")
        p.add_run(" 完成")
        doc.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="split_tpl", format="word", file_path=tpl
        )

        output = str(workspace / "split_out.docx")
        result = doc_handler.doc_apply_template("split_tpl", output, {"title": "年度总结"})
        assert result["variables_replaced"] >= 1

        doc2 = Document(output)
        text = "\n".join(p.text for p in doc2.paragraphs)
        assert "报告：年度总结 完成" in text
        assert "{{" not in text

    def test_apply_ppt_template_split_runs(self, workspace: Path) -> None:
        """PPT 占位符跨 run 时应能替换。"""
        tpl = str(workspace / "split.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tf = tb.text_frame
        tf.text = "标题：{{ti"
        tf.paragraphs[0].add_run().text = "tle}}"
        prs.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="split_ppt", format="ppt", file_path=tpl
        )

        output = str(workspace / "split_ppt_out.pptx")
        result = doc_handler.doc_apply_template("split_ppt", output, {"title": "公司介绍"})
        assert result["variables_replaced"] >= 1

        prs2 = Presentation(output)
        text = prs2.slides[0].shapes[0].text_frame.text
        assert "标题：公司介绍" in text
        assert "{{" not in text

    def test_apply_ppt_template_sections(self, workspace: Path) -> None:
        """doc_apply_template 通过 sections 按章节扩展 PPT 页数。"""
        tpl = workspace / "sec_tpl.pptx"
        _build_ppt_section_template(tpl)
        doc_handler.doc_manage_template(
            action="register", name="sec_tpl", format="ppt", file_path=str(tpl)
        )

        output = str(workspace / "sec_out.pptx")
        result = doc_handler.doc_apply_template(
            "sec_tpl",
            output,
            variables={"title": "公司介绍", "item1": "概况", "item2": "业务", "item3": "优势", "item4": "联系", "contact": "联系"},
            sections=[
                {
                    "section_no": "01", "section_title": "公司概况", "slide_title": "公司概况",
                    "point1": "成立于2019年", "point2": "注册资本500万", "point3": "法人刘欢", "point4": "天府新区",
                },
                {
                    "section_no": "02", "section_title": "主营业务", "slide_title": "主营业务",
                    "point1": "软硬件开发", "point2": "数据处理", "point3": "系统集成", "point4": "咨询",
                },
            ],
        )
        # 封面 + 目录 + 2 章节 x (章节页+内容页) + 结尾 = 7 页
        assert result["slide_count"] == 7
        assert result["template"] == "sec_tpl"

        prs = Presentation(output)
        texts = [prs.slides[i].shapes[0].text_frame.text for i in range(len(prs.slides))]
        assert texts[0] == "公司介绍"
        assert texts[1] == "概况 / 业务 / 优势 / 联系"
        assert texts[2] == "PART 01 公司概况"
        assert "成立于2019年" in texts[3]
        assert "法人刘欢" in texts[3]
        assert texts[4] == "PART 02 主营业务"
        assert "数据处理" in texts[5]
        assert texts[6] == "联系"
        assert not any("{{" in t for t in texts)

    def test_apply_ppt_template_page_numbers(self, workspace: Path) -> None:
        """扩页后应重写克隆页的静态页码，保证最终页码连续。"""
        tpl = str(workspace / "pagenum_tpl.pptx")
        prs = Presentation()
        blank = prs.slide_layouts[6]

        def add_text(slide, text: str, left=1.0, top=1.0) -> None:
            tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(6), Inches(1))
            tb.text_frame.text = text

        def add_footer(slide, text: str) -> None:
            # 底部区域页码（top=7.0in，幻灯片高 7.5in）
            add_text(slide, text, left=11.5, top=7.0)

        s = prs.slides.add_slide(blank)
        add_text(s, "{{title}}")  # 封面：无页码
        s = prs.slides.add_slide(blank)
        add_text(s, "{{item1}} / {{item2}}")
        add_footer(s, "02")  # 目录页
        s = prs.slides.add_slide(blank)
        add_text(s, "PART {{section_no}} {{section_title}}")
        add_footer(s, "03")  # 章节页原型
        s = prs.slides.add_slide(blank)
        add_text(s, "{{slide_title}}\n{{point1}}\n{{point2}}")
        add_footer(s, "04")  # 内容页原型
        s = prs.slides.add_slide(blank)
        add_text(s, "{{contact}}")  # 结尾：无页码
        prs.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="pagenum_tpl", format="ppt", file_path=tpl
        )

        output = str(workspace / "pagenum_out.pptx")
        doc_handler.doc_apply_template(
            "pagenum_tpl",
            output,
            variables={"title": "T", "item1": "A", "item2": "B", "contact": "C"},
            sections=[
                {"section_no": "01", "section_title": "第一章", "slide_title": "第一章", "point1": "p1", "point2": "p2"},
                {"section_no": "02", "section_title": "第二章", "slide_title": "第二章", "point1": "p1", "point2": "p2"},
            ],
        )

        prs2 = Presentation(output)
        assert len(prs2.slides) == 7
        footers = []
        for i in range(len(prs2.slides)):
            nums = [
                sh.text_frame.text.strip()
                for sh in prs2.slides[i].shapes
                if sh.has_text_frame and sh.text_frame.text.strip().isdigit()
            ]
            footers.append(nums[0] if nums else None)
        # 封面/结尾无页码，其余按最终页序 02..06
        assert footers == [None, "02", "03", "04", "05", "06", None]

    def test_apply_ppt_template_sections_without_proto(self, workspace: Path) -> None:
        """模板无章节/内容原型页时，sections 退化为仅做全局填充。"""
        tpl = str(workspace / "plain_tpl.pptx")
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        tb.text_frame.text = "{{name}}"
        prs.save(tpl)
        doc_handler.doc_manage_template(
            action="register", name="plain_tpl", format="ppt", file_path=tpl
        )

        output = str(workspace / "plain_out.pptx")
        result = doc_handler.doc_apply_template(
            "plain_tpl",
            output,
            variables={"name": "银杉"},
            sections=[{"section_no": "01", "section_title": "章节"}],
        )
        assert result["slide_count"] == 1
        prs2 = Presentation(output)
        assert prs2.slides[0].shapes[0].text_frame.text == "银杉"


class TestDocExtractPlaceholders:
    def test_extract(self, template_docx: str) -> None:
        doc_handler.doc_manage_template(
            action="register",
            name="ph_tpl",
            format="word",
            file_path=template_docx,
        )
        result = doc_handler.doc_extract_placeholders("ph_tpl")
        assert result["count"] >= 2
        names = [p["name"] for p in result["placeholders"]]
        assert "title" in names
        assert "author" in names


class TestDocSession:
    def test_open_save_close_session(self, workspace: Path) -> None:
        # 创建文档
        docx_path = str(workspace / "session_test.docx")
        word_handler.word_create_document(docx_path, title="Session 测试")

        # 打开 Session
        result = doc_handler.doc_open_session(docx_path, format="word")
        session_id = result["session_id"]
        assert result["format"] == "word"

        # 在 Session 中编辑
        word_handler.word_add_paragraph(docx_path, text="Session 段落", session_id=session_id)

        # 列出 Session
        listing = doc_handler.doc_list_sessions()
        assert listing["count"] >= 1

        # 保存并关闭 Session
        doc_handler.doc_close_session(session_id, save=True, output_path=docx_path)

        # 验证 Session 已关闭
        listing = doc_handler.doc_list_sessions()
        assert listing["count"] == 0

        # 验证内容已保存
        doc2 = Document(docx_path)
        text = "\n".join(p.text for p in doc2.paragraphs)
        assert "Session 段落" in text

    def test_close_without_save(self, workspace: Path) -> None:
        docx_path = str(workspace / "nosave.docx")
        word_handler.word_create_document(docx_path, title="不保存测试")
        result = doc_handler.doc_open_session(docx_path, format="word")
        session_id = result["session_id"]
        doc_handler.doc_close_session(session_id, save=False)
        listing = doc_handler.doc_list_sessions()
        assert listing["count"] == 0
