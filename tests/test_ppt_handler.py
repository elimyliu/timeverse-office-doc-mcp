"""测试 PPT handler 工具。"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from timeverse_office_doc_mcp.common.error_handler import ToolError
from timeverse_office_doc_mcp.handlers import ppt_handler


@pytest.fixture
def workspace(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> Path:
    """创建临时工作空间并更新 PathGuard 白名单。"""
    ws = tmp_path / "workspace"
    ws.mkdir()
    from timeverse_office_doc_mcp.common.path_guard import path_guard

    monkeypatch.setattr(path_guard, "allowed_dirs", [tmp_path.resolve()])
    return ws


@pytest.fixture
def pptx_path(workspace: Path) -> str:
    return str(workspace / "test.pptx")


class TestPptCreatePresentation:
    def test_create_empty(self, pptx_path: str) -> None:
        result = ppt_handler.ppt_create_presentation(pptx_path)
        assert result["filename"] == str(Path(pptx_path).resolve())
        assert Path(pptx_path).exists()

    def test_create_and_read_back(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        prs = Presentation(pptx_path)
        assert len(prs.slides) >= 0


class TestPptSlideOperations:
    def test_add_slide(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        result = ppt_handler.ppt_add_slide(pptx_path, layout=6, title="测试")
        assert result["slide_idx"] >= 0

    def test_add_and_delete_slide(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        info_before = ppt_handler.ppt_get_info(pptx_path)
        count_before = info_before["slide_count"]

        ppt_handler.ppt_delete_slide(pptx_path, 0)
        info_after = ppt_handler.ppt_get_info(pptx_path)
        assert info_after["slide_count"] == count_before - 1

    def test_add_slide_invalid_layout(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        with pytest.raises(ToolError):
            ppt_handler.ppt_add_slide(pptx_path, layout=99)


class TestPptAddText:
    def test_add_text(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        result = ppt_handler.ppt_add_text(pptx_path, 0, text="你好世界", font_size=24, bold=True)
        assert result["text"] == "你好世界"

    def test_add_text_invalid_slide(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        with pytest.raises(ToolError):
            ppt_handler.ppt_add_text(pptx_path, 99, text="测试")


class TestPptAddTable:
    def test_add_table_with_data(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        data = [["姓名", "年龄"], ["张三", "25"], ["李四", "30"]]
        result = ppt_handler.ppt_add_table(pptx_path, 0, rows=3, cols=2, data=data)
        assert result["rows"] == 3
        assert result["cols"] == 2


class TestPptSetBackground:
    def test_set_background(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        result = ppt_handler.ppt_set_background(pptx_path, 0, color="FF0000")
        assert result["background_color"] == "FF0000"


class TestPptSetSlideNotes:
    def test_set_and_get_notes(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        ppt_handler.ppt_set_slide_notes(pptx_path, 0, notes_text="这是备注")
        result = ppt_handler.ppt_get_slide_notes(pptx_path, 0)
        assert "这是备注" in result["notes"]


class TestPptExtractText:
    def test_extract(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        ppt_handler.ppt_add_text(pptx_path, 0, text="幻灯片文本内容")
        result = ppt_handler.ppt_extract_text(pptx_path)
        assert len(result["slides"]) >= 1
        found = any("幻灯片文本内容" in text for text in result["slides"][0]["texts"])
        assert found


class TestPptAnalyzeStructure:
    def test_analyze(self, pptx_path: str) -> None:
        ppt_handler.ppt_create_presentation(pptx_path)
        ppt_handler.ppt_add_slide(pptx_path, layout=6)
        ppt_handler.ppt_add_text(pptx_path, 0, text="标题")
        ppt_handler.ppt_add_table(pptx_path, 0, rows=2, cols=2, data=[["A", "B"], ["1", "2"]])
        result = ppt_handler.ppt_analyze_structure(pptx_path)
        assert result["slide_count"] >= 1
        assert result["total_shapes"] >= 2
